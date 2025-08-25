"""
Content extraction for various document types.

Handles text extraction from different file formats including
PDF, Word, PowerPoint, HTML, and other document types.
"""

import io
import os
import tempfile
from pathlib import Path
from typing import Dict, Any, Optional, Tuple

from ....shared import get_logger
from ..models import DocumentType


class ContentExtractor:
    """
    Extracts text content from various document formats.
    
    Supports multiple file types and provides metadata extraction
    for enhanced document processing.
    """
    
    def __init__(self):
        self.logger = get_logger(__name__)
        self.logger.info("Content Extractor initialized")
    
    def extract_from_file(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """
        Extract text content from a file.
        
        Args:
            file_path: Path to the file to process
            
        Returns:
            Tuple of (extracted_text, metadata)
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        # Determine document type from extension
        doc_type = self._get_document_type(file_path)
        
        # Extract content based on type
        if doc_type == DocumentType.TEXT:
            return self._extract_text(file_path)
        elif doc_type == DocumentType.MARKDOWN:
            return self._extract_markdown(file_path)
        elif doc_type == DocumentType.HTML:
            return self._extract_html(file_path)
        elif doc_type == DocumentType.PDF:
            return self._extract_pdf(file_path)
        elif doc_type == DocumentType.WORD:
            return self._extract_word(file_path)
        elif doc_type == DocumentType.POWERPOINT:
            return self._extract_powerpoint(file_path)
        elif doc_type == DocumentType.CSV:
            return self._extract_csv(file_path)
        elif doc_type == DocumentType.JSON:
            return self._extract_json(file_path)
        else:
            # Default to plain text extraction
            return self._extract_text(file_path)
    
    def extract_from_content(self, content: str, doc_type: DocumentType) -> Tuple[str, Dict[str, Any]]:
        """
        Extract and process content that's already in string format.
        
        Args:
            content: Raw content string
            doc_type: Type of content
            
        Returns:
            Tuple of (processed_text, metadata)
        """
        metadata = {
            'document_type': doc_type.value,
            'character_count': len(content),
            'word_count': len(content.split()),
            'line_count': len(content.splitlines())
        }
        
        if doc_type == DocumentType.HTML:
            content = self._clean_html_content(content)
        elif doc_type == DocumentType.MARKDOWN:
            content = self._clean_markdown_content(content)
        elif doc_type == DocumentType.JSON:
            content = self._extract_json_text(content)
        
        # Basic text cleaning
        content = self._clean_text(content)
        
        # Update metadata after cleaning
        metadata['processed_character_count'] = len(content)
        metadata['processed_word_count'] = len(content.split())
        
        return content, metadata
    
    def _get_document_type(self, file_path: Path) -> DocumentType:
        """Determine document type from file extension."""
        extension = file_path.suffix.lower()
        
        type_mapping = {
            '.txt': DocumentType.TEXT,
            '.md': DocumentType.MARKDOWN,
            '.markdown': DocumentType.MARKDOWN,
            '.html': DocumentType.HTML,
            '.htm': DocumentType.HTML,
            '.pdf': DocumentType.PDF,
            '.doc': DocumentType.WORD,
            '.docx': DocumentType.WORD,
            '.ppt': DocumentType.POWERPOINT,
            '.pptx': DocumentType.POWERPOINT,
            '.csv': DocumentType.CSV,
            '.json': DocumentType.JSON,
        }
        
        return type_mapping.get(extension, DocumentType.TEXT)
    
    def _extract_text(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Extract content from plain text file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            metadata = {
                'file_size': file_path.stat().st_size,
                'encoding': 'utf-8',
                'line_count': len(content.splitlines())
            }
            
            return self._clean_text(content), metadata
            
        except UnicodeDecodeError:
            # Try with different encoding
            with open(file_path, 'r', encoding='latin-1') as f:
                content = f.read()
            
            metadata = {
                'file_size': file_path.stat().st_size,
                'encoding': 'latin-1',
                'line_count': len(content.splitlines())
            }
            
            return self._clean_text(content), metadata
    
    def _extract_markdown(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Extract content from Markdown file."""
        content, metadata = self._extract_text(file_path)
        
        # Clean markdown formatting while preserving structure
        content = self._clean_markdown_content(content)
        
        metadata['document_type'] = 'markdown'
        return content, metadata
    
    def _extract_html(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Extract content from HTML file."""
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            self.logger.warning("BeautifulSoup not available, using basic HTML extraction")
            return self._extract_text(file_path)
        
        with open(file_path, 'r', encoding='utf-8') as f:
            html_content = f.read()
        
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Extract text content
        content = soup.get_text()
        
        # Extract metadata
        metadata = {
            'file_size': file_path.stat().st_size,
            'title': soup.title.string if soup.title else None,
            'document_type': 'html'
        }
        
        return self._clean_text(content), metadata
    
    def _extract_pdf(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Extract content from PDF file."""
        try:
            import PyPDF2
        except ImportError:
            self.logger.warning("PyPDF2 not available, cannot extract PDF content")
            raise ImportError("PyPDF2 required for PDF extraction")
        
        content = ""
        metadata = {
            'file_size': file_path.stat().st_size,
            'document_type': 'pdf',
            'page_count': 0
        }
        
        try:
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                metadata['page_count'] = len(pdf_reader.pages)
                
                for page in pdf_reader.pages:
                    content += page.extract_text() + "\n"
            
            return self._clean_text(content), metadata
            
        except Exception as e:
            self.logger.error(f"Failed to extract PDF content: {e}")
            raise
    
    def _extract_word(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Extract content from Word document."""
        try:
            from docx import Document as WordDocument
        except ImportError:
            self.logger.warning("python-docx not available, cannot extract Word content")
            raise ImportError("python-docx required for Word extraction")
        
        try:
            doc = WordDocument(file_path)
            
            content = ""
            for paragraph in doc.paragraphs:
                content += paragraph.text + "\n"
            
            metadata = {
                'file_size': file_path.stat().st_size,
                'document_type': 'word',
                'paragraph_count': len(doc.paragraphs)
            }
            
            return self._clean_text(content), metadata
            
        except Exception as e:
            self.logger.error(f"Failed to extract Word content: {e}")
            raise
    
    def _extract_powerpoint(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Extract content from PowerPoint presentation."""
        try:
            from pptx import Presentation
        except ImportError:
            self.logger.warning("python-pptx not available, cannot extract PowerPoint content")
            raise ImportError("python-pptx required for PowerPoint extraction")
        
        try:
            prs = Presentation(file_path)
            
            content = ""
            slide_count = 0
            
            for slide in prs.slides:
                slide_count += 1
                content += f"--- Slide {slide_count} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
                
                content += "\n"
            
            metadata = {
                'file_size': file_path.stat().st_size,
                'document_type': 'powerpoint',
                'slide_count': slide_count
            }
            
            return self._clean_text(content), metadata
            
        except Exception as e:
            self.logger.error(f"Failed to extract PowerPoint content: {e}")
            raise
    
    def _extract_csv(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Extract content from CSV file."""
        try:
            import pandas as pd
        except ImportError:
            # Fallback to basic CSV processing
            return self._extract_csv_basic(file_path)
        
        try:
            df = pd.read_csv(file_path)
            
            # Convert to text representation
            content = f"CSV Data:\n"
            content += f"Columns: {', '.join(df.columns)}\n\n"
            content += df.to_string(index=False)
            
            metadata = {
                'file_size': file_path.stat().st_size,
                'document_type': 'csv',
                'row_count': len(df),
                'column_count': len(df.columns),
                'columns': list(df.columns)
            }
            
            return content, metadata
            
        except Exception as e:
            self.logger.error(f"Failed to extract CSV content: {e}")
            return self._extract_csv_basic(file_path)
    
    def _extract_csv_basic(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Basic CSV extraction without pandas."""
        import csv
        
        content = "CSV Data:\n"
        row_count = 0
        columns = []
        
        with open(file_path, 'r', encoding='utf-8') as f:
            csv_reader = csv.reader(f)
            
            for i, row in enumerate(csv_reader):
                if i == 0:
                    columns = row
                    content += f"Columns: {', '.join(row)}\n\n"
                else:
                    content += f"Row {i}: {', '.join(row)}\n"
                    row_count += 1
        
        metadata = {
            'file_size': file_path.stat().st_size,
            'document_type': 'csv',
            'row_count': row_count,
            'column_count': len(columns),
            'columns': columns
        }
        
        return content, metadata
    
    def _extract_json(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Extract content from JSON file."""
        import json
        
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        content = self._extract_json_text(json.dumps(data, indent=2))
        
        metadata = {
            'file_size': file_path.stat().st_size,
            'document_type': 'json'
        }
        
        return content, metadata
    
    def _extract_json_text(self, json_str: str) -> str:
        """Convert JSON to readable text format."""
        import json
        
        try:
            data = json.loads(json_str) if isinstance(json_str, str) else json_str
            
            def flatten_json(obj, parent_key='', sep='_'):
                items = []
                if isinstance(obj, dict):
                    for k, v in obj.items():
                        new_key = f"{parent_key}{sep}{k}" if parent_key else k
                        if isinstance(v, (dict, list)):
                            items.extend(flatten_json(v, new_key, sep).items())
                        else:
                            items.append((new_key, str(v)))
                elif isinstance(obj, list):
                    for i, v in enumerate(obj):
                        new_key = f"{parent_key}{sep}{i}" if parent_key else str(i)
                        if isinstance(v, (dict, list)):
                            items.extend(flatten_json(v, new_key, sep).items())
                        else:
                            items.append((new_key, str(v)))
                return dict(items)
            
            flattened = flatten_json(data)
            content = "\n".join([f"{k}: {v}" for k, v in flattened.items()])
            return content
            
        except Exception as e:
            self.logger.error(f"Failed to process JSON content: {e}")
            return str(json_str)
    
    def _clean_html_content(self, html_content: str) -> str:
        """Clean HTML content to extract plain text."""
        try:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(html_content, 'html.parser')
            return soup.get_text()
        except ImportError:
            # Basic HTML cleaning without BeautifulSoup
            import re
            # Remove HTML tags
            content = re.sub(r'<[^>]+>', '', html_content)
            return content
    
    def _clean_markdown_content(self, markdown_content: str) -> str:
        """Clean Markdown formatting while preserving structure."""
        import re
        
        # Remove markdown formatting but keep structure
        content = markdown_content
        
        # Remove emphasis markers but keep text
        content = re.sub(r'\*\*(.*?)\*\*', r'\1', content)  # Bold
        content = re.sub(r'\*(.*?)\*', r'\1', content)      # Italic
        content = re.sub(r'~~(.*?)~~', r'\1', content)      # Strikethrough
        
        # Remove link formatting but keep text
        content = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', content)
        
        # Remove code formatting
        content = re.sub(r'`([^`]+)`', r'\1', content)
        content = re.sub(r'```[^`]*```', '', content)
        
        # Clean up headers but preserve the text
        content = re.sub(r'^#{1,6}\s*(.*)$', r'\1', content, flags=re.MULTILINE)
        
        return content
    
    def _clean_text(self, text: str) -> str:
        """Basic text cleaning and normalization."""
        if not text:
            return ""
        
        # Normalize whitespace
        import re
        
        # Replace multiple whitespace with single space
        text = re.sub(r'\s+', ' ', text)
        
        # Remove excessive newlines
        text = re.sub(r'\n\s*\n\s*\n+', '\n\n', text)
        
        # Strip leading/trailing whitespace
        text = text.strip()
        
        return text